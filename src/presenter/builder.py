"""
Step 3: Presentation Builder Module.
Populates TECMOTIV_SOLUTIONS.pptx template preserving slide master branding,
applying Tecmotiv typography, card layouts, and speaker notes.
"""

import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.config import (
    DEFAULT_TEMPLATE_PATH,
    COLOR_PRIMARY_NAVY,
    COLOR_ACCENT_TEAL,
    COLOR_VIOLET_PRIMARY,
    COLOR_CARD_BG,
    COLOR_TEXT_MAIN,
    COLOR_TEXT_MUTED,
    COLOR_WHITE,
    FONT_HEADING,
    FONT_BODY,
)
from src.models import PresentationContent, SlideContent

logger = logging.getLogger("PresentationBuilder")


class PresentationBuilder:
    """Builds PPTX presentation populated with generated content based on Tecmotiv template."""

    def __init__(self, template_path: str = DEFAULT_TEMPLATE_PATH):
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template file not found at path: {template_path}")
        self.template_path = template_path

    def build_presentation(self, content: PresentationContent, output_pptx_path: str) -> str:
        """
        Clones template and builds all slides off the master layout.
        """
        logger.info(f"Starting Step 3: Populating template '{self.template_path}' with {len(content.slides)} slides...")

        prs = Presentation(self.template_path)
        blank_layout = prs.slide_layouts[0]

        # Dynamically update bottom footer banner date in Slide Master
        from datetime import datetime
        current_month_year = datetime.now().strftime("%B %Y")
        current_year = datetime.now().strftime("%Y")

        try:
            for shape in blank_layout.slide_master.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if "June 2026" in run.text:
                                run.text = run.text.replace("June 2026", current_month_year)
                            elif "2026" in run.text and current_year != "2026":
                                run.text = run.text.replace("2026", current_year)
        except Exception as e:
            logger.warning(f"Could not update footer banner date in master layout: {e}")

        # Delete existing initial blank slides in template to build fresh content set
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        # Build each slide
        for index, slide_data in enumerate(content.slides):
            slide = prs.slides.add_slide(blank_layout)
            if index == 0:
                self._build_title_slide(slide, slide_data, content.topic)
            else:
                self._build_content_slide(slide, slide_data)

            # Add speaker notes safely
            if slide_data.speaker_notes:
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide and notes_slide.notes_text_frame:
                        notes_slide.notes_text_frame.text = slide_data.speaker_notes
                except Exception as e:
                    logger.warning(f"Could not attach speaker notes to slide {index + 1}: {e}")

        try:
            if os.path.exists(output_pptx_path):
                try:
                    os.remove(output_pptx_path)
                except Exception:
                    pass
            prs.save(output_pptx_path)
        except PermissionError:
            alt_path = output_pptx_path.replace(".pptx", "_generated.pptx")
            prs.save(alt_path)
            output_pptx_path = alt_path

        logger.info(f"Step 3 Complete: Presentation saved to '{output_pptx_path}'")
        return output_pptx_path

    def _build_title_slide(self, slide, data: SlideContent, topic: str):
        """Constructs Slide 1 (Title Slide) with executive presentation hero branding."""
        # Top Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.4), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = data.subtitle_tag.upper() if data.subtitle_tag else "STRATEGIC PRESENTATION"
        p_tag.alignment = PP_ALIGN.LEFT
        p_tag.font.name = FONT_BODY
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_ACCENT_TEAL

        # Main Deck Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.4))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = data.title
        p_title.alignment = PP_ALIGN.LEFT
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_NAVY

        # Subtle Accent Horizontal Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(2.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT_TEAL
        line.line.color.rgb = COLOR_ACCENT_TEAL

        # Executive Context Card Container
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.9), Inches(8.4), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_ACCENT_TEAL
        card.line.width = Pt(1)

        # Executive Metadata inside Title Card
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.3)
        tf_card.margin_top = Inches(0.2)
        tf_card.margin_right = Inches(0.3)

        from datetime import datetime
        current_date_str = datetime.now().strftime("%B %Y")

        title_metadata = [
            "Presented by: Tecmotiv Strategic Advisory",
            f"Date: {current_date_str}",
            f"Strategic Objective: {topic} Executive Discussion Document"
        ]

        for i, line_text in enumerate(title_metadata):
            p = tf_card.paragraphs[0] if i == 0 else tf_card.add_paragraph()
            p.text = line_text
            p.alignment = PP_ALIGN.LEFT
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_TEXT_MAIN
            p.space_after = Pt(8)

    def _build_content_slide(self, slide, data: SlideContent):
        """Constructs Content Slide with category header, content card, and bottom takeaway bar."""
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.4), Inches(0.3))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        tf_tag.margin_top = Inches(0)
        tf_tag.margin_bottom = Inches(0)
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = data.subtitle_tag.upper()
        p_tag.alignment = PP_ALIGN.LEFT
        p_tag.font.name = FONT_BODY
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_ACCENT_TEAL

        # Dynamic layout spacing based on title length
        is_long_title = len(data.title) > 42
        title_font_size = Pt(18) if is_long_title else Pt(22)
        card_top = Inches(1.48) if is_long_title else Inches(1.35)
        card_height = Inches(2.62) if is_long_title else Inches(2.75)

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.62), Inches(8.4), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_top = Inches(0)
        tf_title.margin_bottom = Inches(0)
        p_title = tf_title.paragraphs[0]
        p_title.text = data.title
        p_title.alignment = PP_ALIGN.LEFT
        p_title.font.name = FONT_HEADING
        p_title.font.size = title_font_size
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_NAVY

        # Main Content Card Background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_top, Inches(8.4), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_ACCENT_TEAL
        card.line.width = Pt(0.75)

        # Bullet Points Text inside Card
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.3)
        tf_card.margin_top = Inches(0.18)
        tf_card.margin_right = Inches(0.3)
        tf_card.margin_bottom = Inches(0.18)

        bullet_font_size = Pt(11.5) if is_long_title or len(data.bullet_points) >= 4 else Pt(13)
        bullet_spacing = Pt(4.5) if is_long_title or len(data.bullet_points) >= 4 else Pt(7)

        for i, bullet in enumerate(data.bullet_points):
            p = tf_card.paragraphs[0] if i == 0 else tf_card.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = bullet_spacing

            raw_bullet = bullet.strip()
            
            # Case A: Bullet has explicit markdown bold **Header:** format
            if raw_bullet.startswith("**") and "**" in raw_bullet[2:]:
                parts = raw_bullet[2:].split("**", 1)
                lead_in = parts[0].rstrip(":").strip()
                desc = parts[1].lstrip(":").strip()
            # Case B: Bullet has colon-separated Lead-In Header
            elif ":" in raw_bullet and len(raw_bullet.split(":", 1)[0]) <= 55:
                lead_in, desc = raw_bullet.split(":", 1)
                lead_in = lead_in.replace("**", "").strip()
                desc = desc.replace("**", "").strip()
            else:
                lead_in = None
                desc = raw_bullet.replace("**", "").strip()

            if lead_in:
                # Run 1: Bullet dot + Bold Lead-In Header
                r1 = p.add_run()
                r1.text = f"•  {lead_in}: "
                r1.font.name = FONT_BODY
                r1.font.size = bullet_font_size
                r1.font.bold = True
                r1.font.color.rgb = COLOR_PRIMARY_NAVY

                # Run 2: Explanation body text
                r2 = p.add_run()
                r2.text = desc
                r2.font.name = FONT_BODY
                r2.font.size = bullet_font_size
                r2.font.bold = False
                r2.font.color.rgb = COLOR_TEXT_MAIN
            else:
                r = p.add_run()
                r.text = f"•  {desc}"
                r.font.name = FONT_BODY
                r.font.size = bullet_font_size
                r.font.bold = False
                r.font.color.rgb = COLOR_TEXT_MAIN

        # Key Takeaway Banner at bottom of content area
        if data.key_takeaway:
            takeaway_top = Inches(4.2)
            takeaway_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), takeaway_top, Inches(8.4), Inches(0.65)
            )
            takeaway_box.fill.solid()
            takeaway_box.fill.fore_color.rgb = COLOR_PRIMARY_NAVY
            takeaway_box.line.color.rgb = COLOR_ACCENT_TEAL
            takeaway_box.line.width = Pt(1.5)

            tf_takeaway = takeaway_box.text_frame
            tf_takeaway.word_wrap = True
            tf_takeaway.margin_left = Inches(0.2)
            tf_takeaway.margin_right = Inches(0.2)
            tf_takeaway.margin_top = Inches(0.08)
            tf_takeaway.margin_bottom = Inches(0.08)
            p_takeaway = tf_takeaway.paragraphs[0]
            p_takeaway.text = f"KEY TAKEAWAY: {data.key_takeaway}"
            p_takeaway.font.name = FONT_BODY
            p_takeaway.font.size = Pt(10.5)
            p_takeaway.font.bold = True
            p_takeaway.font.color.rgb = COLOR_WHITE
            p_takeaway.alignment = PP_ALIGN.LEFT
